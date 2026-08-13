import os
import shutil
import logging
import fitz  # PyMuPDF
import cv2
import numpy as np
from PIL import Image
from app.core.config import settings

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".docx", ".pptx", ".ppt"}
MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB


def deskew_image(img: np.ndarray) -> np.ndarray:
    """Deskew image using text contour minimum area rectangle"""
    try:
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) if len(img.shape) == 3 else img
        thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)[1]
        coords = np.column_stack(np.where(thresh > 0))
        
        if len(coords) < 10:
            return img

        angle = cv2.minAreaRect(coords)[-1]
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle

        if abs(angle) > 0.5 and abs(angle) < 45.0:
            (h, w) = img.shape[:2]
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, angle, 1.0)
            img = cv2.warpAffine(img, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
            logger.info(f"Deskewed image by {angle:.2f} degrees")
    except Exception as e:
        logger.warning(f"Deskew failed, returning original image: {str(e)}")
    return img


def enhance_contrast(img: np.ndarray) -> np.ndarray:
    """Apply CLAHE contrast enhancement and light Gaussian smoothing"""
    try:
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img

        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(gray)
        smoothed = cv2.GaussianBlur(enhanced, (3, 3), 0)
        return cv2.cvtColor(smoothed, cv2.COLOR_GRAY2BGR)
    except Exception as e:
        logger.warning(f"Enhance contrast failed, returning original image: {str(e)}")
        return img


def clean_and_save_image(image_path: str, output_path: str) -> None:
    """Read image, perform deskew & contrast enhancement, write to output_path"""
    img = cv2.imread(image_path)
    if img is None:
        # Fallback to Pillow if OpenCV fails to read directly
        pil_img = Image.open(image_path)
        img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)

    img = deskew_image(img)
    img = enhance_contrast(img)
    cv2.imwrite(output_path, img)


def render_pdf_pages_to_images(pdf_path: str, processed_dir: str) -> int:
    """Render PDF pages into high-res PNG images using PyMuPDF"""
    os.makedirs(processed_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    page_count = len(doc)

    for i in range(page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=300)
        raw_page_path = os.path.join(processed_dir, f"raw_page_{i+1}.png")
        final_page_path = os.path.join(processed_dir, f"page_{i+1}.png")
        
        pix.save(raw_page_path)
        clean_and_save_image(raw_page_path, final_page_path)
        
        if os.path.exists(raw_page_path):
            os.remove(raw_page_path)

    doc.close()
    return page_count


def render_docx_to_preview_image(docx_path: str, output_path: str) -> None:
    """Render DOCX text and table structure onto a high-quality document preview canvas"""
    img = np.ones((1200, 850, 3), dtype=np.uint8) * 255  # Clean white paper canvas
    
    # Blue top accent bar for document header
    cv2.rectangle(img, (0, 0), (850, 40), (220, 150, 60), -1)
    cv2.putText(img, "DOCX DOCUMENT PREVIEW", (30, 26), cv2.FONT_HERSHEY_SIMPLEX, 0.6, (255, 255, 255), 2)
    
    y = 80
    try:
        import docx
        doc = docx.Document(docx_path)
        
        # Render Paragraphs
        for p in doc.paragraphs:
            p_text = p.text.strip()
            if not p_text:
                continue
            
            font_scale = 0.55 if y < 140 else 0.42
            color = (0, 0, 0)
            thickness = 2 if y < 140 else 1
            
            words = p_text.split()
            current_line = ""
            for word in words:
                test_line = f"{current_line} {word}".strip()
                if len(test_line) * 9 > 780:
                    cv2.putText(img, current_line[:90], (40, y), cv2.FONT_HERSHEY_SIMPLEX, font_scale, color, thickness)
                    y += 24
                    current_line = word
                else:
                    current_line = test_line
            if current_line:
                cv2.putText(img, current_line[:90], (40, y), cv2.FONT_HERSHEY_SIMPLEX, font_scale, color, thickness)
                y += 26
            
            if y > 1100:
                break
                
        # Render Tables
        for t_idx, table in enumerate(doc.tables):
            if y > 1050:
                break
            y += 10
            cv2.putText(img, f"--- Table {t_idx + 1} ---", (40, y), cv2.FONT_HERSHEY_SIMPLEX, 0.45, (180, 100, 0), 1)
            y += 20
            
            for row in table.rows[:10]:
                row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                row_str = " | ".join(row_cells)[:95]
                cv2.putText(img, row_str, (45, y), cv2.FONT_HERSHEY_SIMPLEX, 0.38, (40, 40, 40), 1)
                y += 20
                if y > 1100:
                    break
    except Exception as e:
        logger.warning(f"DOCX image preview render error: {str(e)}")
        cv2.putText(img, "DOCX Document Content (Text Stream Active)", (40, y), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 1)

    cv2.imwrite(output_path, img)


def render_pptx_to_preview_image(pptx_path: str, output_path: str) -> int:
    """Render PPTX slide previews onto document preview canvas"""
    img = np.ones((900, 1200, 3), dtype=np.uint8) * 255  # 16:9 presentation canvas
    
    # Purple top accent header
    cv2.rectangle(img, (0, 0), (1200, 50), (180, 80, 120), -1)
    cv2.putText(img, "POWERPOINT PRESENTATION PREVIEW", (30, 34), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (255, 255, 255), 2)
    
    slide_count = 1
    y = 100
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        for idx, slide in enumerate(prs.slides[:5], start=1):
            cv2.putText(img, f"Slide {idx} / {slide_count}", (40, y), cv2.FONT_HERSHEY_SIMPLEX, 0.55, (180, 80, 120), 2)
            y += 26
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        first_line = text.splitlines()[0][:110]
                        cv2.putText(img, f"  • {first_line}", (50, y), cv2.FONT_HERSHEY_SIMPLEX, 0.45, (30, 30, 30), 1)
                        y += 24
                        if y > 820:
                            break
            y += 15
            if y > 820:
                break
    except Exception as e:
        logger.warning(f"PPTX preview render error: {str(e)}")
        cv2.putText(img, "PPTX Presentation Slides (Text Stream Active)", (40, y), cv2.FONT_HERSHEY_SIMPLEX, 0.6, (0, 0, 0), 1)

    cv2.imwrite(output_path, img)
    return slide_count


def preprocess_document(document_id: str, file_path: str, storage_dir: str = None) -> dict:
    """Main preprocessing pipeline entrypoint"""
    base_storage = storage_dir or settings.STORAGE_DIR
    doc_dir = os.path.join(base_storage, document_id)
    processed_dir = os.path.join(doc_dir, "processed")
    os.makedirs(processed_dir, exist_ok=True)

    ext = os.path.splitext(file_path)[1].lower()
    page_count = 1

    if ext == ".pdf":
        page_count = render_pdf_pages_to_images(file_path, processed_dir)
    elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".tif"}:
        final_page_path = os.path.join(processed_dir, "page_1.png")
        clean_and_save_image(file_path, final_page_path)
        page_count = 1
    elif ext == ".docx":
        final_page_path = os.path.join(processed_dir, "page_1.png")
        shutil.copyfile(file_path, os.path.join(processed_dir, "original.docx"))
        render_docx_to_preview_image(file_path, final_page_path)
        page_count = 1
    elif ext in {".pptx", ".ppt"}:
        final_page_path = os.path.join(processed_dir, "page_1.png")
        shutil.copyfile(file_path, os.path.join(processed_dir, f"original{ext}"))
        page_count = render_pptx_to_preview_image(file_path, final_page_path)

    return {
        "document_id": document_id,
        "page_count": page_count,
        "processed_dir": processed_dir
    }
