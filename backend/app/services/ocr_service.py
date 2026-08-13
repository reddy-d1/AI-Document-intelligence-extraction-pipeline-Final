import os
import json
import logging
from typing import Dict, Any, List, Optional
import fitz  # PyMuPDF

from app.core.config import settings
from app.services.ocr.base import DocumentOCRResult, PageOCRResult, WordBoundingBox
from app.services.ocr.factory import get_ocr_provider

logger = logging.getLogger(__name__)


def extract_native_pdf_text(pdf_path: str) -> Optional[DocumentOCRResult]:
    """Attempt direct vector text extraction from native text-based PDFs via PyMuPDF"""
    try:
        doc = fitz.open(pdf_path)
        page_results: List[PageOCRResult] = []
        full_text_list: List[str] = []
        total_words = 0

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text").strip()
            words_data = page.get_text("words")  # (x0, y0, x1, y1, word, block_no, line_no, word_no)
            
            word_boxes: List[WordBoundingBox] = []
            for w in words_data:
                word_boxes.append(
                    WordBoundingBox(
                        word=w[4],
                        x=int(w[0]),
                        y=int(w[1]),
                        w=int(w[2] - w[0]),
                        h=int(w[3] - w[1]),
                        confidence=1.0,
                        page_num=page_num + 1
                    )
                )

            total_words += len(word_boxes)
            full_text_list.append(text)
            page_results.append(
                PageOCRResult(
                    page_num=page_num + 1,
                    text=text,
                    confidence=1.0,
                    word_boxes=word_boxes
                )
            )

        doc.close()

        # If PDF has native text stream (>= 3 words total)
        if total_words >= 3:
            raw_text = "\n\n--- Page Break ---\n\n".join(full_text_list)

            return DocumentOCRResult(
                raw_text=raw_text,
                confidence=1.0,
                pages=page_results,
                extraction_method="native_pdf"
            )
    except Exception as e:
        logger.warning(f"Native PDF text extraction failed for {pdf_path}: {str(e)}")

    return None


def extract_native_docx_text(docx_path: str) -> Optional[DocumentOCRResult]:
    """Extract paragraphs and structured tables from native machine-readable DOCX files"""
    try:
        import docx
        doc = docx.Document(docx_path)
        text_parts: List[str] = []
        word_boxes: List[WordBoundingBox] = []
        word_counter = 0

        # 1. Paragraphs
        for p in doc.paragraphs:
            p_text = p.text.strip()
            if p_text:
                text_parts.append(p_text)
                for w in p_text.split():
                    word_counter += 1
                    word_boxes.append(
                        WordBoundingBox(
                            word=w,
                            x=50,
                            y=word_counter * 18,
                            w=len(w) * 8,
                            h=14,
                            confidence=1.0,
                            page_num=1
                        )
                    )

        # 2. Structured Tables (preserving row & column relationships)
        for t_idx, table in enumerate(doc.tables):
            table_lines: List[str] = []
            table_lines.append(f"\n--- Table {t_idx + 1} ---")
            for r_idx, row in enumerate(table.rows):
                row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                row_str = " | ".join(row_cells)
                table_lines.append(f"| {row_str} |")
                
                for cell_text in row_cells:
                    for w in cell_text.split():
                        word_counter += 1
                        word_boxes.append(
                            WordBoundingBox(
                                word=w,
                                x=50,
                                y=word_counter * 18,
                                w=len(w) * 8,
                                h=14,
                                confidence=1.0,
                                page_num=1
                            )
                        )
            text_parts.append("\n".join(table_lines))

        raw_text = "\n\n".join(text_parts)

        page_result = PageOCRResult(
            page_num=1,
            text=raw_text,
            confidence=1.0,
            word_boxes=word_boxes
        )

        return DocumentOCRResult(
            raw_text=raw_text,
            confidence=1.0,
            pages=[page_result],
            extraction_method="native_docx"
        )
    except Exception as e:
        logger.warning(f"Native DOCX text extraction failed for {docx_path}: {str(e)}")

    return None


def extract_native_pptx_text(pptx_path: str) -> Optional[DocumentOCRResult]:
    """Extract slide text, shapes, bullet points, tables, and speaker notes from PPTX presentations"""
    try:
        page_results: List[PageOCRResult] = []
        text_parts: List[str] = []
        word_boxes: List[WordBoundingBox] = []
        word_counter = 0

        # Try python-pptx first
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_lines: List[str] = []
                slide_lines.append(f"--- Slide {slide_num} ---")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            p_text = paragraph.text.strip()
                            if p_text:
                                slide_lines.append(p_text)
                                for w in p_text.split():
                                    word_counter += 1
                                    word_boxes.append(WordBoundingBox(
                                        word=w, x=50, y=word_counter * 18, w=len(w) * 8, h=14, confidence=1.0, page_num=slide_num
                                    ))
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_str = " | ".join([c.text.strip().replace("\n", " ") for c in row.cells])
                            slide_lines.append(f"| {row_str} |")

                # Check for Speaker Notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"[Speaker Notes: {notes}]")

                slide_text = "\n".join(slide_lines)
                text_parts.append(slide_text)
                page_results.append(PageOCRResult(
                    page_num=slide_num,
                    text=slide_text,
                    confidence=1.0,
                    word_boxes=word_boxes
                ))

            raw_text = "\n\n".join(text_parts)
            return DocumentOCRResult(
                raw_text=raw_text,
                confidence=1.0,
                pages=page_results,
                extraction_method="native_pptx"
            )

        except Exception as e_pptx:
            logger.warning(f"python-pptx extraction fallback notice: {str(e_pptx)}. Using native OpenXML zip parser.")

        # OpenXML Fallback via standard zipfile & xml parser
        import zipfile
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            slide_files = sorted([f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])
            if not slide_files:
                return None

            for page_idx, slide_file in enumerate(slide_files, start=1):
                slide_xml = zip_ref.read(slide_file)
                tree = ET.fromstring(slide_xml)
                texts = [node.text for node in tree.iter() if node.tag.endswith('}t') and node.text]
                slide_text = " ".join(texts).strip()

                full_slide_content = f"--- Slide {page_idx} ---\n{slide_text}"
                text_parts.append(full_slide_content)
                page_results.append(PageOCRResult(
                    page_num=page_idx, text=full_slide_content, confidence=1.0, word_boxes=[]
                ))

            raw_text = "\n\n".join(text_parts)
            return DocumentOCRResult(
                raw_text=raw_text, confidence=1.0, pages=page_results, extraction_method="native_pptx_xml"
            )
    except Exception as e:
        logger.warning(f"Native PPTX text extraction failed for {pptx_path}: {str(e)}")

    return None


def process_document_ocr(document_id: str, file_path: str, storage_dir: str = None) -> DocumentOCRResult:
    """Main OCR pipeline function for documents"""
    base_storage = storage_dir or settings.STORAGE_DIR
    doc_dir = os.path.join(base_storage, document_id)
    processed_dir = os.path.join(doc_dir, "processed")
    ext = os.path.splitext(file_path)[1].lower()

    # 1. Try native PDF text stream extraction first
    if ext == ".pdf":
        native_res = extract_native_pdf_text(file_path)
        if native_res:
            logger.info(f"Extracted native vector text for PDF {document_id}")
            save_ocr_metadata(doc_dir, native_res)
            return native_res
    elif ext == ".docx":
        docx_res = extract_native_docx_text(file_path)
        if docx_res:
            logger.info(f"Extracted native paragraph & table text for DOCX {document_id}")
            save_ocr_metadata(doc_dir, docx_res)
            return docx_res
    elif ext in {".pptx", ".ppt"}:
        pptx_res = extract_native_pptx_text(file_path)
        if pptx_res:
            logger.info(f"Extracted native slide & speaker note text for PPTX {document_id}")
            save_ocr_metadata(doc_dir, pptx_res)
            return pptx_res

    # 2. Fall back to image OCR (Tesseract / Cloud) on preprocessed page images
    ocr_provider = get_ocr_provider()
    page_results: List[PageOCRResult] = []
    full_text_list: List[str] = []
    confidences: List[float] = []

    # Find preprocessed page images page_1.png, page_2.png, etc.
    page_num = 1
    while True:
        page_img = os.path.join(processed_dir, f"page_{page_num}.png")
        if not os.path.exists(page_img):
            break

        logger.info(f"Running OCR on page {page_num} for document {document_id}")
        page_res = ocr_provider.extract_page_text(page_img, page_num=page_num)
        page_results.append(page_res)
        full_text_list.append(page_res.text)
        confidences.append(page_res.confidence)
        page_num += 1

    # If no preprocessed images found, attempt directly on file_path
    if not page_results and os.path.exists(file_path):
        page_res = ocr_provider.extract_page_text(file_path, page_num=1)
        page_results.append(page_res)
        full_text_list.append(page_res.text)
        confidences.append(page_res.confidence)

    raw_text = "\n\n--- Page Break ---\n\n".join(full_text_list)
    avg_confidence = round(sum(confidences) / len(confidences), 4) if confidences else 0.0

    ocr_result = DocumentOCRResult(
        raw_text=raw_text,
        confidence=avg_confidence,
        pages=page_results,
        extraction_method=settings.OCR_PROVIDER
    )

    save_ocr_metadata(doc_dir, ocr_result)
    return ocr_result


def save_ocr_metadata(doc_dir: str, ocr_result: DocumentOCRResult) -> None:
    """Save full OCR bounding box JSON metadata to /storage/{document_id}/ocr_data.json"""
    ocr_json_path = os.path.join(doc_dir, "ocr_data.json")
    with open(ocr_json_path, "w", encoding="utf-8") as f:
        json.dump(ocr_result.model_dump(), f, indent=2, ensure_ascii=False)
    logger.info(f"Saved OCR metadata to {ocr_json_path}")
